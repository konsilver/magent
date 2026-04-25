"""Type-aware summary builder for uploaded artifacts.

Used to generate a short text snippet that is injected into the agent's
prompt as "historical file" context, so the agent knows what it can later
fetch via the `read_artifact` tool.

Each generator is best-effort: it must not raise, and must return a short
string (typically under 800 chars). Any internal failure falls back to the
generic summary.
"""

from __future__ import annotations

import csv
import io
import logging
import os
from typing import Optional

logger = logging.getLogger(__name__)


# Total cap for any generated summary — defense-in-depth in case a generator misbehaves.
_MAX_SUMMARY_CHARS = 1000


def build_summary_from_text(text: str, filename: str, mime_type: str = "") -> str:
    """Derive a summary from already-parsed text — no bytes-level re-parse.

    Prefer this path over `build_summary(bytes, ...)` whenever the caller
    already has the parsed text (e.g., the frontend parsed the file and
    sent it as `attachment.content`, or `/v1/file/parse` was called).

    For types with no meaningful text (images), returns a name+mime summary.
    """
    if not text or not text.strip():
        return _summary_generic_from_meta(filename, mime_type)

    text = text.strip()
    ext = os.path.splitext(filename.lower())[1]
    mime = (mime_type or "").lower()

    try:
        # Line-oriented types — preserve structure (headings / code lines)
        if ext == ".md" or ext in _CODE_EXTS:
            lines = text.splitlines()
            head = "\n".join(lines[:30])
            return _cap(f"{head}\n（共 {len(lines)} 行）")

        # CSV parsed to markdown table — show first 6 rows
        if ext == ".csv":
            lines = text.splitlines()
            head = "\n".join(lines[:6])
            return _cap(f"{head}\n（共 {len(lines)} 行）")

        # XLSX parsed to markdown with `### Sheet: N` headers — keep headers
        # so the model sees sheet names + first rows.  Wider cap since the
        # structural hints are valuable.
        if ext in {".xlsx", ".xls"} or mime.endswith("spreadsheetml.sheet"):
            return _cap(_truncate_preserving_sheet_headers(text, 800))

        # Default: first 500 chars + total length hint
        return _cap(text[:500].rstrip() + f"\n（共 {len(text):,} 字）")
    except Exception as e:
        logger.warning(f"build_summary_from_text failed for {filename}: {e}")
        return _summary_generic_from_meta(filename, mime_type)


def _truncate_preserving_sheet_headers(text: str, budget: int) -> str:
    """Cut XLSX/markdown text at sheet boundaries when possible."""
    if len(text) <= budget:
        return text
    # Find last sheet boundary within budget
    snippet = text[:budget]
    last_sheet = snippet.rfind("### Sheet")
    if last_sheet > 0:
        # Prefer cutting at a line boundary AFTER the last visible sheet
        tail = text[last_sheet:]
        # Include enough context: next 400 chars from that sheet header
        cut = min(len(tail), max(400, budget - last_sheet))
        return text[:last_sheet] + tail[:cut] + "\n…（更多 sheet/行 未展示）"
    return snippet.rstrip() + "\n…"


def _summary_generic_from_meta(filename: str, mime_type: str) -> str:
    return f"{filename} · {mime_type or '未知类型'}"


def _cap(s: str) -> str:
    if len(s) > _MAX_SUMMARY_CHARS:
        return s[:_MAX_SUMMARY_CHARS] + "…"
    return s


def build_summary(file_bytes: bytes, filename: str, mime_type: str = "") -> str:
    """Dispatch to the type-specific summary generator.

    Returns a short string describing the file's contents. Never raises —
    falls back to a generic summary on any internal error.
    """
    try:
        ext = os.path.splitext(filename.lower())[1]
        mime = (mime_type or "").lower()

        if mime.startswith("image/") or ext in {".png", ".jpg", ".jpeg", ".gif", ".webp", ".bmp"}:
            out = _summary_image(file_bytes, filename, mime)
        elif ext in {".xlsx", ".xls"}:
            out = _summary_xlsx(file_bytes, filename)
        elif ext == ".csv":
            out = _summary_csv(file_bytes, filename)
        elif ext == ".pptx":
            out = _summary_pptx(file_bytes, filename)
        elif ext == ".md":
            out = _summary_markdown(file_bytes, filename)
        elif ext in _CODE_EXTS:
            out = _summary_code(file_bytes, filename)
        elif ext in {".pdf", ".docx", ".doc", ".wps", ".txt"}:
            out = _summary_text_document(file_bytes, filename)
        else:
            out = _summary_generic(file_bytes, filename, mime)
    except Exception as e:
        logger.warning(f"build_summary failed for {filename}: {e}")
        out = _summary_generic(file_bytes, filename, mime_type)

    if len(out) > _MAX_SUMMARY_CHARS:
        out = out[:_MAX_SUMMARY_CHARS] + "…"
    return out


# ── Text documents (PDF/DOCX/DOC/WPS/TXT) ─────────────────────────────────────

def _summary_text_document(file_bytes: bytes, filename: str) -> str:
    """Parse and truncate first 500 chars."""
    from core.content.file_parser import parse_file

    text = parse_file(file_bytes, filename) or ""
    if not text:
        return _summary_generic(file_bytes, filename, "")

    head = text[:500].rstrip()
    total = len(text)
    suffix = f"\n（共 {total:,} 字）" if total > 500 else f"\n（共 {total} 字）"
    return head + suffix


# ── Markdown ──────────────────────────────────────────────────────────────────

def _summary_markdown(file_bytes: bytes, filename: str) -> str:
    """First 30 lines of markdown, preserving heading structure."""
    text = _decode_bytes(file_bytes)
    lines = text.splitlines()
    head = "\n".join(lines[:30])
    total_lines = len(lines)
    return f"{head}\n（共 {total_lines} 行）"


# ── Code files ────────────────────────────────────────────────────────────────

_CODE_EXTS = {
    ".py", ".js", ".ts", ".tsx", ".jsx", ".java", ".go", ".rs", ".c", ".cpp",
    ".h", ".hpp", ".cs", ".rb", ".php", ".sh", ".sql", ".yaml", ".yml",
    ".json", ".xml", ".html", ".css", ".vue",
}


def _summary_code(file_bytes: bytes, filename: str) -> str:
    """First 30 lines of code."""
    text = _decode_bytes(file_bytes)
    lines = text.splitlines()
    head = "\n".join(lines[:30])
    total_lines = len(lines)
    return f"{head}\n（共 {total_lines} 行）"


# ── XLSX / XLS (use openpyxl directly — faster than parse_file for summary) ──

def _summary_xlsx(file_bytes: bytes, filename: str) -> str:
    """For XLSX: list sheet names, headers, first 3 rows per sheet.

    Uses openpyxl directly to avoid parsing the entire workbook to markdown.
    Falls back to parse_file() for .xls or when openpyxl fails.
    """
    ext = os.path.splitext(filename.lower())[1]
    if ext == ".xls":
        # xlrd path is slower but consistent; just use parse_file's output.
        return _summary_text_document(file_bytes, filename)

    try:
        from openpyxl import load_workbook
    except ImportError:
        return _summary_text_document(file_bytes, filename)

    try:
        wb = load_workbook(io.BytesIO(file_bytes), read_only=True, data_only=True)
        sheets = wb.sheetnames
        parts: list[str] = [f"Excel 工作簿（共 {len(sheets)} 个 sheet）："]
        total_len = len(parts[0])

        for sheet_name in sheets:
            if total_len > 800:
                parts.append(f"（还有 {len(sheets) - (len(parts) - 1)} 个 sheet 未展示）")
                break
            ws = wb[sheet_name]
            sheet_part = _summarize_one_sheet(sheet_name, ws)
            parts.append(sheet_part)
            total_len += len(sheet_part)

        wb.close()
        return "\n".join(parts)
    except Exception as e:
        logger.warning(f"openpyxl summary failed for {filename}: {e}")
        return _summary_text_document(file_bytes, filename)


def _summarize_one_sheet(sheet_name: str, ws) -> str:
    """Build a compact summary of one worksheet: headers + first 3 rows."""
    rows_iter = ws.iter_rows(max_row=4, values_only=True)
    rows = [r for r in rows_iter if r is not None]
    if not rows:
        return f"- Sheet「{sheet_name}」: 空"

    def _cell(v):
        if v is None:
            return ""
        s = str(v).strip()
        return s if len(s) <= 30 else s[:30] + "…"

    header = " | ".join(_cell(v) for v in rows[0])
    data_lines = []
    for r in rows[1:4]:
        data_lines.append(" | ".join(_cell(v) for v in r))

    lines = [f"- Sheet「{sheet_name}」"]
    if header:
        lines.append(f"  表头: {header}")
    for dl in data_lines:
        if dl.strip(" |"):
            lines.append(f"  数据: {dl}")
    return "\n".join(lines)


# ── CSV ───────────────────────────────────────────────────────────────────────

def _summary_csv(file_bytes: bytes, filename: str) -> str:
    """Column names + first 5 rows."""
    text = _decode_bytes(file_bytes)
    reader = csv.reader(io.StringIO(text))
    rows: list[list[str]] = []
    for i, row in enumerate(reader):
        rows.append(row)
        if i >= 5:
            break

    if not rows:
        return f"CSV 文件「{filename}」: 空"

    header = " | ".join(rows[0][:10])  # cap columns
    parts = [f"CSV 文件「{filename}」", f"列名: {header}"]
    for r in rows[1:]:
        parts.append("数据: " + " | ".join(c[:30] for c in r[:10]))
    return "\n".join(parts)


# ── PPTX ──────────────────────────────────────────────────────────────────────

def _summary_pptx(file_bytes: bytes, filename: str) -> str:
    """Slide count + titles of first 20 slides."""
    try:
        from pptx import Presentation  # type: ignore
    except ImportError:
        return _summary_generic(file_bytes, filename, "application/vnd.openxmlformats-officedocument.presentationml.presentation")

    try:
        prs = Presentation(io.BytesIO(file_bytes))
        slides = list(prs.slides)
        titles: list[str] = []
        for i, slide in enumerate(slides[:20]):
            title = ""
            try:
                if slide.shapes.title is not None:
                    title = (slide.shapes.title.text or "").strip()
            except Exception:
                pass
            if not title:
                title = "（无标题）"
            titles.append(f"  第 {i+1} 页: {title[:60]}")

        parts = [f"PPTX 演示文稿（共 {len(slides)} 页）:"]
        parts.extend(titles)
        if len(slides) > 20:
            parts.append(f"  （后续 {len(slides) - 20} 页未展示）")
        return "\n".join(parts)
    except Exception as e:
        logger.warning(f"pptx summary failed for {filename}: {e}")
        return _summary_generic(file_bytes, filename, "application/vnd.openxmlformats-officedocument.presentationml.presentation")


# ── Image ─────────────────────────────────────────────────────────────────────

def _summary_image(file_bytes: bytes, filename: str, mime: str) -> str:
    """Filename + mime + dimensions if PIL available."""
    size_kb = max(1, len(file_bytes) // 1024)
    mime_label = mime or "image/*"

    try:
        from PIL import Image  # type: ignore
        img = Image.open(io.BytesIO(file_bytes))
        w, h = img.size
        return f"[图片] {filename} · {mime_label} · {w}×{h}px · {size_kb}KB"
    except Exception:
        return f"[图片] {filename} · {mime_label} · {size_kb}KB"


# ── Generic fallback ──────────────────────────────────────────────────────────

def _summary_generic(file_bytes: bytes, filename: str, mime: str) -> str:
    size_kb = max(1, len(file_bytes) // 1024)
    return f"{filename} · {mime or '未知类型'} · {size_kb}KB"


# ── Helpers ───────────────────────────────────────────────────────────────────

def _decode_bytes(file_bytes: bytes) -> str:
    """Try UTF-8 then GBK, finally replace errors."""
    for enc in ("utf-8", "gbk", "utf-16"):
        try:
            return file_bytes.decode(enc)
        except UnicodeDecodeError:
            continue
    return file_bytes.decode("utf-8", errors="replace")
